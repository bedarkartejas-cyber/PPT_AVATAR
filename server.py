import os
import random
import logging
import json
import comtypes.client  # Client to talk to PowerPoint
import pythoncom        # <--- REQUIRED FOR THREADING FIX
from flask import Flask, jsonify, request, send_from_directory
from flask_cors import CORS
from dotenv import load_dotenv
from livekit import api
from werkzeug.utils import secure_filename
from pptx import Presentation

load_dotenv()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger('ppt-server')

app = Flask(__name__, static_folder='.', static_url_path='')
CORS(app)

# Folders for storage
UPLOAD_FOLDER = 'uploads'
SLIDES_FOLDER = 'slides'  
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SLIDES_FOLDER, exist_ok=True)

# --- 1. PPT PROCESSOR (FIXED FOR THREADING) ---
def process_ppt(ppt_path):
    slides_data = []
    
    # *** CRITICAL FIX: Initialize Windows COM for this thread ***
    pythoncom.CoInitialize() 
    
    # A. EXPORT IMAGES (Robust Method)
    try:
        abs_ppt_path = os.path.abspath(ppt_path)
        abs_slides_folder = os.path.abspath(SLIDES_FOLDER)

        # Clear old slides
        for f in os.listdir(abs_slides_folder):
            file_path = os.path.join(abs_slides_folder, f)
            if os.path.isfile(file_path):
                os.remove(file_path)

        print(f"⏳ Connecting to PowerPoint...")
        # Force a new instance to avoid conflicts
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        # powerpoint.Visible = 1 # Uncomment if you want to see it happen
        
        presentation = powerpoint.Presentations.Open(abs_ppt_path, WithWindow=False)
        
        print(f"📸 Exporting {len(presentation.Slides)} slides...")
        
        for i, slide in enumerate(presentation.Slides, start=1):
            image_path = os.path.join(abs_slides_folder, f"Slide{i}.jpg")
            slide.Export(image_path, "JPG")
            print(f"   -> Saved Slide{i}.jpg")

        presentation.Close()
        # powerpoint.Quit() 
        
        print(f"✅ All images exported to {SLIDES_FOLDER}/")

    except Exception as e:
        print(f"❌ PowerPoint Error: {e}")
        # Note: We don't return [] here, we continue to extract text 
        # so the agent can still read even if images fail.

    # B. EXTRACT TEXT (For the Agent)
    try:
        if not os.path.exists(ppt_path):
            return []

        prs = Presentation(ppt_path)
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            # Get Title
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Title: {slide.shapes.title.text}")

            # Get Other Text
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            image_filename = f"Slide{i}.jpg"
            
            # Verify if image actually exists
            if not os.path.exists(os.path.join(SLIDES_FOLDER, image_filename)):
                print(f"⚠️ Warning: Image for Slide {i} was not found.")

            slides_data.append({
                "slide_number": i,
                "image_url": f"/slides/{image_filename}", 
                "content": " ".join(slide_text)
            })

        # Save data for Agent
        with open("presentation.json", "w", encoding="utf-8") as f:
            json.dump(slides_data, f, indent=4, ensure_ascii=False)

        return slides_data

    except Exception as e:
        logger.error(f"❌ Error reading PPT text: {e}")
        return []
    finally:
        # Good practice to uninitialize
        pythoncom.CoUninitialize()

# --- 2. ROUTES ---

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')

@app.route('/slides/<path:filename>')
def serve_slide(filename):
    return send_from_directory(SLIDES_FOLDER, filename)

@app.route('/api/upload-ppt', methods=['POST'])
def upload_ppt():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file and file.filename.endswith('.pptx'):
        filename = secure_filename(file.filename)
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)
        
        # Process the file
        data = process_ppt(filepath)
        
        if not data:
             return jsonify({"error": "Processing failed. Check server logs."}), 500

        return jsonify({
            "status": "success", 
            "message": "Presentation processed",
            "slide_count": len(data)
        })
    
    return jsonify({"error": "Invalid file type"}), 400

@app.route('/api/connection-details')
def connection_details():
    try:
        LIVEKIT_URL = os.getenv('LIVEKIT_URL')
        API_KEY = os.getenv('LIVEKIT_API_KEY')
        API_SECRET = os.getenv('LIVEKIT_API_SECRET')

        if not all([LIVEKIT_URL, API_KEY, API_SECRET]):
            return jsonify({"error": "Missing Keys"}), 500

        room_name = f"ppt_session_{random.randint(10000, 99999)}"
        participant_identity = f"user_{random.randint(1000, 9999)}"
        
        token = api.AccessToken(API_KEY, API_SECRET) \
            .with_identity(participant_identity) \
            .with_name("User") \
            .with_grants(api.VideoGrants(
                room_join=True,
                room=room_name,
                can_publish=True,
                can_subscribe=True,
            )).to_jwt()

        return jsonify({
            "serverUrl": LIVEKIT_URL,
            "roomName": room_name,
            "participantToken": token
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

def start_server():
    from waitress import serve
    print("🟢 Web Server running on http://localhost:8000")
    serve(app, host='0.0.0.0', port=8000, threads=4)

if __name__ == '__main__':
    start_server()